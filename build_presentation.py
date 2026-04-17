"""
Build ARIA presentation for executives and compliance teams.
Output: framl_ai_presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
from pathlib import Path

# ── AI Image Generation ───────────────────────────────────────────────────────
IMG_DIR = Path("docs/img_gen")
IMG_DIR.mkdir(parents=True, exist_ok=True)

_c = lambda r, g, b: (r/255, g/255, b/255)
_DARK = _c(0x1B, 0x2A, 0x4A)
_MID  = _c(0x1F, 0x5C, 0x99)
_ORG  = _c(0xF0, 0x8C, 0x00)
_LGHT = _c(0xD6, 0xE4, 0xF0)
_CARD = _c(0xE8, 0xF1, 0xF8)


def make_neural_net(path, w=5.5, h=3.1):
    """Glowing neural network on dark background — for title slide."""
    fig, ax = plt.subplots(figsize=(w, h))
    fig.patch.set_facecolor(_DARK)
    ax.set_facecolor(_DARK)
    ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.axis('off')
    np.random.seed(42)
    layers = [3, 5, 6, 5, 2]
    xs = np.linspace(0.1, 0.9, len(layers))
    nodes = [list(zip([x]*n, np.linspace(0.12, 0.88, n))) for x, n in zip(xs, layers)]
    for i in range(len(layers) - 1):
        for x1, y1 in nodes[i]:
            for x2, y2 in nodes[i + 1]:
                a = np.random.uniform(0.05, 0.28)
                col = _ORG if np.random.random() > 0.78 else _MID
                ax.plot([x1, x2], [y1, y2], color=col, alpha=a, lw=0.7, zorder=1)
    for i, layer in enumerate(nodes):
        for x, y in layer:
            nc = _ORG if i in (0, len(nodes) - 1) else _LGHT
            for sz, al in [(900, 0.06), (400, 0.15), (120, 0.5), (30, 1.0)]:
                c = _MID if sz > 300 else nc
                ax.scatter(x, y, s=sz, c=[c], alpha=al, zorder=3)
    fig.savefig(str(path), dpi=150, bbox_inches='tight', pad_inches=0.0)
    plt.close(fig)


def make_cluster_scatter(path, w=3.4, h=1.35):
    """4-cluster scatter — for Segmentation demo card."""
    np.random.seed(7)
    fig, ax = plt.subplots(figsize=(w, h))
    fig.patch.set_facecolor(_LGHT); ax.set_facecolor(_LGHT); ax.axis('off')
    centers = [(0.22, 0.62), (0.55, 0.78), (0.45, 0.28), (0.78, 0.52)]
    colors  = [_DARK, _MID, _ORG, _c(0x2E, 0x86, 0xAB)]
    for (cx, cy), col in zip(centers, colors):
        xs = np.random.normal(cx, 0.065, 28); ys = np.random.normal(cy, 0.065, 28)
        ax.scatter(xs, ys, c=[col], alpha=0.65, s=14, zorder=2)
        ax.scatter(cx, cy, c=[col], s=90, marker='*', zorder=3,
                   edgecolors='white', linewidths=0.5)
    ax.set_xlim(0, 1); ax.set_ylim(0, 1)
    fig.savefig(str(path), dpi=120, bbox_inches='tight', pad_inches=0.06)
    plt.close(fig)


def make_fp_fn_curve(path, w=3.4, h=1.35):
    """FP/FN trade-off crossing curves — for Threshold demo card."""
    fig, ax = plt.subplots(figsize=(w, h))
    fig.patch.set_facecolor(_LGHT); ax.set_facecolor(_LGHT)
    t = np.linspace(0, 1, 200)
    fp = 1 - t**0.55
    fn = t**1.3
    ax.plot(t, fp, color=_ORG, lw=2.2, label='FP Rate')
    ax.plot(t, fn, color=_MID, lw=2.2, label='FN Rate')
    cross = np.argmin(np.abs(fp - fn))
    ax.axvline(t[cross], color=_DARK, lw=1.4, ls='--', alpha=0.7)
    ax.scatter([t[cross]], [fp[cross]], color=_DARK, s=55, zorder=5)
    ax.legend(fontsize=7, loc='center right', framealpha=0.85,
              edgecolor=_MID, facecolor=_LGHT)
    ax.set_xlabel('Threshold →', fontsize=7, color=_DARK)
    ax.set_xticks([]); ax.set_yticks([])
    for sp in ax.spines.values():
        sp.set_color(_MID); sp.set_linewidth(0.8)
    fig.tight_layout(pad=0.3)
    fig.savefig(str(path), dpi=120, bbox_inches='tight', pad_inches=0.06)
    plt.close(fig)


def make_rag_pipeline(path, w=3.4, h=1.35):
    """Docs → Embed → VectorDB → Answer pipeline — for Policy demo card."""
    fig, ax = plt.subplots(figsize=(w, h))
    fig.patch.set_facecolor(_LGHT); ax.set_facecolor(_LGHT)
    ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.axis('off')
    stages = [('Docs', 0.13), ('Embed', 0.38), ('Vector\nDB', 0.63), ('Answer', 0.88)]
    for label, x in stages:
        circle = plt.Circle((x, 0.52), 0.11, color=_MID, zorder=2)
        ax.add_patch(circle)
        ax.text(x, 0.52, label, ha='center', va='center', fontsize=7,
                color='white', fontweight='bold', zorder=3)
    for i in range(len(stages) - 1):
        x1 = stages[i][1] + 0.11; x2 = stages[i+1][1] - 0.11
        ax.annotate('', xy=(x2, 0.52), xytext=(x1, 0.52),
                    arrowprops=dict(arrowstyle='->', color=_ORG, lw=1.8))
    fig.savefig(str(path), dpi=120, bbox_inches='tight', pad_inches=0.06)
    plt.close(fig)


def make_lora_diagram(path, w=4.0, h=1.65):
    """Frozen layers + LoRA adapter blocks — for Fine-Tuning slide."""
    fig, ax = plt.subplots(figsize=(w, h))
    fig.patch.set_facecolor(_CARD); ax.set_facecolor(_CARD)
    ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.axis('off')
    ys = np.linspace(0.82, 0.18, 5)
    bh = 0.12
    for i, y in enumerate(ys):
        ax.add_patch(plt.Rectangle((0.04, y - bh/2), 0.54, bh,
                                   color=_MID, alpha=0.88, zorder=2))
        ax.text(0.31, y, 'Frozen', ha='center', va='center',
                fontsize=6.5, color='white', fontweight='bold', zorder=3)
        if i in (1, 3):
            for bx, lbl in [(0.63, 'A'), (0.80, 'B')]:
                ax.add_patch(plt.Rectangle((bx, y - bh/2), 0.14, bh,
                                           color=_ORG, zorder=2))
                ax.text(bx + 0.07, y, lbl, ha='center', va='center',
                        fontsize=6.5, color='white', fontweight='bold', zorder=3)
            ax.plot([0.58, 0.63], [y, y], color=_ORG, lw=1.2, zorder=1)
            ax.plot([0.77, 0.80], [y, y], color=_ORG, lw=1.2, zorder=1)
    ax.text(0.31, 0.04, '← Frozen (base model)', ha='center', va='bottom',
            fontsize=5.5, color=_MID)
    ax.text(0.755, 0.04, '← LoRA adapters\n(trainable ~1%)', ha='center', va='bottom',
            fontsize=5.5, color=_ORG)
    fig.savefig(str(path), dpi=120, bbox_inches='tight', pad_inches=0.05)
    plt.close(fig)


print("Generating AI images...")
IMG_NEURAL  = IMG_DIR / "neural_net.png"
IMG_CLUSTER = IMG_DIR / "cluster_scatter.png"
IMG_FPFN    = IMG_DIR / "fp_fn_curve.png"
IMG_RAG     = IMG_DIR / "rag_pipeline.png"
IMG_LORA    = IMG_DIR / "lora_diagram.png"
make_neural_net(IMG_NEURAL)
make_cluster_scatter(IMG_CLUSTER)
make_fp_fn_curve(IMG_FPFN)
make_rag_pipeline(IMG_RAG)
make_lora_diagram(IMG_LORA)
print("Images ready.")

# ── Colours ──────────────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1B, 0x2A, 0x4A)   # header bars / title slide
MID_BLUE    = RGBColor(0x1F, 0x5C, 0x99)   # card titles, accents
LIGHT_BLUE  = RGBColor(0xD6, 0xE4, 0xF0)   # card backgrounds
CARD_BLUE   = RGBColor(0xE8, 0xF1, 0xF8)   # very light card fill
ORANGE      = RGBColor(0xF0, 0x8C, 0x00)   # highlights / accent strip
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GREY   = RGBColor(0x40, 0x40, 0x40)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank

# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, l, t, w, h, text, size=18, bold=False, color=DARK_GREY,
                align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def add_bullet_box(slide, l, t, w, h, title, bullets, title_size=16, bullet_size=13,
                   title_color=MID_BLUE, bullet_color=DARK_GREY, bg=None, indent=True):
    if bg:
        add_rect(slide, l, t, w, h, fill=bg)
    txb = slide.shapes.add_textbox(Inches(l+0.15), Inches(t+0.1),
                                   Inches(w-0.3), Inches(h-0.2))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = Pt(title_size)
    run.font.bold = True
    run.font.color.rgb = title_color

    for b in bullets:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if indent:
            p.level = 1
        run = p.add_run()
        run.text = f"• {b}"
        run.font.size = Pt(bullet_size)
        run.font.color.rgb = bullet_color

def header_bar(slide, title, subtitle=None):
    """Dark blue top bar with title."""
    add_rect(slide, 0, 0, 13.33, 1.2, fill=DARK_BLUE)
    add_textbox(slide, 0.3, 0.15, 12.0, 0.65, title,
                size=28, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, 0.3, 0.78, 12.0, 0.38, subtitle,
                    size=14, bold=False, color=LIGHT_BLUE)
    # Orange accent line
    add_rect(slide, 0, 1.18, 13.33, 0.06, fill=ORANGE)

_slide_counter = [0]

def footer(slide, text="Princeton Strategy Group  |  Confidential"):
    _slide_counter[0] += 1
    add_rect(slide, 0, 7.2, 13.33, 0.3, fill=DARK_BLUE)
    add_textbox(slide, 0.3, 7.22, 10.5, 0.25, text,
                size=9, color=LIGHT_BLUE, align=PP_ALIGN.LEFT)
    add_textbox(slide, 11.8, 7.22, 1.3, 0.25, str(_slide_counter[0]),
                size=9, color=LIGHT_BLUE, align=PP_ALIGN.RIGHT)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title slide
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill=DARK_BLUE)
add_rect(slide, 0, 3.3, 13.33, 0.08, fill=ORANGE)
add_textbox(slide, 1.0, 1.5, 11.0, 1.2,
            "ARIA",
            size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, 1.0, 2.8, 11.0, 0.7,
            "Dynamic Segmentation  ·  Threshold Tuning  ·  AML Policy Q&A",
            size=20, bold=False, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_textbox(slide, 1.0, 3.6, 11.0, 0.5,
            "Executive & Compliance Overview",
            size=16, color=ORANGE, align=PP_ALIGN.CENTER)
add_textbox(slide, 1.0, 6.8, 11.0, 0.4,
            "March 2026  |  Confidential",
            size=11, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
slide.shapes.add_picture(str(IMG_NEURAL), Inches(7.4), Inches(3.85), Inches(5.6), Inches(3.1))

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — System Architecture Diagram
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(slide, "System Architecture", "Component overview and data flow")
footer(slide)

CX = 13.33 / 2  # 6.665

# ── User Browser ──
add_rect(slide, CX-1.3, 1.35, 2.6, 0.42, fill=LIGHT_BLUE, line=MID_BLUE)
add_textbox(slide, CX-1.3, 1.37, 2.6, 0.38, "User  /  Browser",
            size=13, bold=True, color=MID_BLUE, align=PP_ALIGN.CENTER)
add_textbox(slide, CX-0.15, 1.79, 0.3, 0.22, "↓", size=13, bold=True,
            color=ORANGE, align=PP_ALIGN.CENTER)

# ── Dash Web Application ──
add_rect(slide, 1.5, 2.03, 10.33, 0.5, fill=LIGHT_BLUE, line=MID_BLUE)
add_textbox(slide, 1.5, 2.05, 10.33, 0.46,
            "Dash Web Application   (Flask  ·  Chat UI  ·  Side-panel quick prompts)",
            size=14, bold=True, color=MID_BLUE, align=PP_ALIGN.CENTER)
add_textbox(slide, CX-0.15, 2.55, 0.3, 0.22, "↓", size=13, bold=True,
            color=ORANGE, align=PP_ALIGN.CENTER)

# ── Orchestrator ──
add_rect(slide, CX-2.5, 2.79, 5.0, 0.52, fill=MID_BLUE)
add_textbox(slide, CX-2.5, 2.81, 5.0, 0.48,
            "Orchestrator Agent   (intent routing  ·  agent dispatch)",
            size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── Branch lines from Orchestrator to 3 agents ──
# Agent layout: x=0.3 / 4.92 / 9.33, w=3.5 → centers: 2.05 / 6.67 / 11.08
add_rect(slide, 2.05, 3.44, 9.03, 0.05, fill=ORANGE)          # horizontal bar
add_rect(slide, CX-0.025, 3.31, 0.05, 0.13, fill=ORANGE)      # drop from orchestrator
for ax in [2.05, 6.67, 11.08]:                                 # drops to each agent
    add_rect(slide, ax-0.025, 3.44, 0.05, 0.1, fill=ORANGE)

# ── Agent cards ──
agents_info = [
    (0.3,  "Threshold Agent",    "threshold_tuning\nsar_backtest\nsegment_stats"),
    (4.92, "Segmentation Agent", "ss_cluster_analysis\nalerts_distribution"),
    (9.33, "Policy Agent",       "ChromaDB RAG\n11 AML docs  ·  241 chunks"),
]
for ax, name, tools_text in agents_info:
    add_rect(slide, ax, 3.54, 3.5, 1.55, fill=LIGHT_BLUE)
    add_rect(slide, ax, 3.54, 3.5, 0.52, fill=MID_BLUE)
    add_textbox(slide, ax+0.1, 3.56, 3.3, 0.48,
                name, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, ax+0.15, 4.12, 3.2, 0.9,
                tools_text, size=12, color=DARK_GREY)

# ── Arrows: agents → data sources ──
for ax in [2.05, 6.67, 11.08]:
    add_textbox(slide, ax-0.15, 5.11, 0.3, 0.2, "↓", size=12, bold=True,
                color=ORANGE, align=PP_ALIGN.CENTER)

# ── Data sources ──
add_rect(slide, 0.3, 5.32, 8.12, 0.52, fill=LIGHT_GREY, line=MID_BLUE)
add_textbox(slide, 0.3, 5.34, 8.12, 0.48,
            "Alerts CSV   (5,035 rows  ·  customer transaction aggregates  ·  FP / FN labels)",
            size=12, color=DARK_BLUE, align=PP_ALIGN.CENTER)

add_rect(slide, 9.33, 5.32, 3.5, 0.52, fill=LIGHT_GREY, line=MID_BLUE)
add_textbox(slide, 9.33, 5.34, 3.5, 0.48,
            "ChromaDB   (vector store\nall-MiniLM-L6-v2 embeddings)",
            size=11, color=DARK_BLUE, align=PP_ALIGN.CENTER)

# ── Arrows: data → LM ──
for ax in [4.36, 11.08]:
    add_textbox(slide, ax-0.15, 5.86, 0.3, 0.2, "↓", size=12, bold=True,
                color=ORANGE, align=PP_ALIGN.CENTER)

# ── Fine-Tuned LM ──
add_rect(slide, 0.3, 6.07, 12.73, 0.92, fill=DARK_BLUE)
add_textbox(slide, 0.5, 6.12, 12.33, 0.38,
            "Fine-Tuned Language Model  —  qwen-framl-v5   (Q4_K_M GGUF  ·  LM Studio  ·  localhost:1234/v1)",
            size=13, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 6.5, 12.33, 0.38,
            "Qwen 2.5 7B Instruct  ·  LoRA fine-tuned on 189 FRAML examples  ·  8,192 token context  ·  Serves all agents",
            size=12, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — System Overview
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(slide, "System Overview",
           "End-to-end architecture: from user query to analytical insight")
footer(slide)

# Three-layer architecture diagram
layers = [
    ("User Interface", "Dash web app\nChat interface\nSide-panel quick prompts", 0.4),
    ("Orchestrator", "Intent routing\nQuery classification\nMulti-agent coordination", 4.7),
    ("Agents + Tools", "Threshold Agent  →  threshold_tuning / sar_backtest\n"
                       "Segmentation Agent  →  ss_cluster_analysis / alerts_distribution\n"
                       "Policy Agent  →  ChromaDB RAG retrieval", 9.0),
]
for label, desc, left in layers:
    add_rect(slide, left, 1.5, 3.6, 2.8, fill=LIGHT_BLUE)
    add_textbox(slide, left+0.1, 1.6, 3.4, 0.55, label,
                size=15, bold=True, color=MID_BLUE)
    add_textbox(slide, left+0.1, 2.2, 3.4, 2.0, desc,
                size=13, color=DARK_GREY)

# Arrows between boxes
for x in [4.1, 8.7]:
    add_textbox(slide, x, 2.7, 0.5, 0.4, "→", size=22, bold=True,
                color=ORANGE, align=PP_ALIGN.CENTER)

# Bottom row — model
add_rect(slide, 1.5, 4.6, 10.0, 1.5, fill=LIGHT_BLUE)
add_textbox(slide, 1.6, 4.7, 9.8, 0.5,
            "Fine-Tuned Language Model (qwen-framl-v5  ·  Q4_K_M GGUF  ·  LM Studio API)",
            size=13, bold=True, color=MID_BLUE)
add_textbox(slide, 1.6, 5.2, 9.8, 0.8,
            "Qwen 2.5 7B Instruct base  ·  LoRA fine-tuned on 189 FRAML examples  ·  "
            "8,192 token context  ·  localhost:1234/v1",
            size=13, color=DARK_GREY)

add_textbox(slide, 5.8, 4.3, 1.5, 0.35, "↓  API calls", size=11,
            color=DARK_GREY, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Agents
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(slide, "The Three Agents",
           "Specialised AI agents, each with their own tools and system prompt")
footer(slide)

agent_data = [
    ("Threshold Agent",
     ["Analyses FP/FN trade-offs as alert thresholds change",
      "Tools: threshold_tuning, sar_backtest, segment_stats",
      "Returns PRE-COMPUTED sweep tables — model copies verbatim",
      "Covers Business and Individual segments across 3 threshold columns"]),
    ("Segmentation Agent",
     ["K-Means clustering of customer behavioural profiles",
      "Tools: ss_cluster_analysis, alerts_distribution",
      "Identifies high-risk clusters by transaction frequency and volume",
      "Renders treemap and scatter visualisations"]),
    ("Policy Agent",
     ["Retrieval-Augmented Generation over AML policy documents",
      "ChromaDB knowledge base: 11 docs, 241 chunks",
      "Sources: FFIEC Manual, Wolfsberg Guidance, 31 CFR Parts 1010/1020, FinCEN advisories",
      "Cites source documents; disclaims when KB has no relevant content"]),
]

for i, (title, bullets) in enumerate(agent_data):
    left = 0.3 + i * 4.35
    add_rect(slide, left, 1.45, 4.1, 5.5, fill=LIGHT_BLUE)
    # Coloured title banner inside card
    add_rect(slide, left, 1.45, 4.1, 0.65, fill=MID_BLUE)
    add_textbox(slide, left+0.15, 1.52, 3.8, 0.55,
                title, size=15, bold=True, color=WHITE)
    for j, b in enumerate(bullets):
        add_textbox(slide, left+0.2, 2.2 + j*1.1, 3.8, 1.0,
                    f"• {b}", size=13, color=DARK_GREY)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Tool Interaction
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(slide, "Tool Interaction Flow",
           "How the model calls tools and generates accurate, grounded responses")
footer(slide)

steps = [
    ("1  User Query", "\"What threshold gives a 90%\nSAR catch rate for Individual\nmonthly transaction amount?\"", 0.3),
    ("2  Agent selects tool", "Threshold Agent identifies:\nsar_backtest(\n  segment=Individual,\n  threshold_column=\n  TRXN_AMT_MONTHLY\n)", 3.1),
    ("3  Tool executes", "Python computes sweep against\nreal customer data.\nReturns PRE-COMPUTED block\nwith exact numbers.", 5.9),
    ("4  Model responds", "Copies PRE-COMPUTED block\nverbatim. Adds ONE sentence\nof AML insight. Stops.\nNo hallucinated numbers.", 8.7),
]

for label, desc, left in steps:
    add_rect(slide, left, 1.45, 2.7, 4.5, fill=LIGHT_BLUE)
    add_rect(slide, left, 1.45, 2.7, 0.6, fill=MID_BLUE)
    add_textbox(slide, left+0.1, 1.52, 2.5, 0.5,
                label, size=13, bold=True, color=WHITE)
    add_textbox(slide, left+0.1, 2.1, 2.5, 3.7,
                desc, size=13, color=DARK_GREY)

for x in [3.1, 5.9, 8.7]:
    add_textbox(slide, x-0.35, 3.4, 0.5, 0.4, "→",
                size=22, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

add_rect(slide, 0.3, 6.15, 12.7, 0.85, fill=CARD_BLUE)
add_textbox(slide, 0.5, 6.22, 12.3, 0.7,
            "Key design principle: the model never computes numbers — Python does. "
            "The PRE-COMPUTED block is generated by application.py and passed to the model as the tool result. "
            "The model's only job is to copy it verbatim and add one insight sentence.",
            size=12, color=DARK_BLUE)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Knowledge Base
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(slide, "AML Policy Knowledge Base",
           "Retrieval-Augmented Generation over regulatory and compliance documents")
footer(slide)

# Left — document list
add_rect(slide, 0.3, 1.45, 5.8, 5.5, fill=LIGHT_BLUE)
add_rect(slide, 0.3, 1.45, 5.8, 0.65, fill=MID_BLUE)
add_textbox(slide, 0.45, 1.52, 5.5, 0.55,
            "Source Documents", size=15, bold=True, color=WHITE)
docs = [
    "FFIEC BSA/AML Examination Manual",
    "Wolfsberg Risk-Based Approach Guidance (2006, 2015, 2025)",
    "31 CFR Part 1010 — General BSA Provisions",
    "31 CFR Part 1020 — Rules for Banks",
    "FinCEN Advisory FIN-2014-A008",
    "FinCEN Guidance FIN-2010-G004",
    "BSA/AML Risk Assessment Framework",
    "FFIEC Compliance Program Assessment",
    "BsaAmlExamProcsPackage (4 volumes)",
]
for i, d in enumerate(docs):
    add_textbox(slide, 0.55, 2.2 + i*0.52, 5.4, 0.48,
                f"• {d}", size=12, color=DARK_GREY)

# Right — RAG flow
add_rect(slide, 6.4, 1.45, 6.6, 5.5, fill=LIGHT_GREY)
add_rect(slide, 6.4, 1.45, 6.6, 0.65, fill=MID_BLUE)
add_textbox(slide, 6.55, 1.52, 6.3, 0.55,
            "How RAG Works", size=15, bold=True, color=WHITE)

rag_steps = [
    ("Ingest", "PDFs and DOCX chunked and embedded\nusing all-MiniLM-L6-v2 into ChromaDB\n(241 chunks total)"),
    ("Retrieve", "User query embedded → top-8 semantically\nsimilar chunks retrieved from ChromaDB"),
    ("Generate", "Retrieved chunks injected into model prompt.\nModel cites source document by name.\nIf no relevant content found → disclaims\nand provides general guidance only."),
]
for i, (step, desc) in enumerate(rag_steps):
    add_rect(slide, 6.5, 2.2 + i*1.55, 6.3, 1.45, fill=CARD_BLUE)
    add_textbox(slide, 6.65, 2.25 + i*1.55, 6.0, 0.4,
                step, size=13, bold=True, color=MID_BLUE)
    add_textbox(slide, 6.65, 2.62 + i*1.55, 6.0, 0.95,
                desc, size=13, color=DARK_GREY)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Training Data Pipeline
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(slide, "Training Data Pipeline",
           "How real customer data becomes fine-tuning examples")
footer(slide)

pipeline_steps = [
    ("Step 1\nCustomer Aggregation", "1,090,815 account rows\n↓\n280,621 unique customers\nTransaction metrics summed\nacross all accounts per customer"),
    ("Step 2\nSAR Simulation", "4,835 alerted customers\n10% SAR filing rate = 484 SARs\nSAR score: logistic model on\nalert amount, count, volume"),
    ("Step 3\nFN Injection", "200 non-alerted customers\nadded as false negatives\n(flew under current threshold)\nSource: 20th–80th percentile\nof alerted population"),
    ("Step 4\nAlerts CSV", "5,035 rows\nFP=4,351 | FN=684\nIndividual=4,668 | Business=367\nUsed by all threshold\nand segmentation tools"),
]

for i, (title, desc) in enumerate(pipeline_steps):
    left = 0.3 + i * 3.2
    add_rect(slide, left, 1.45, 3.0, 4.0, fill=LIGHT_BLUE)
    add_rect(slide, left, 1.45, 3.0, 0.75, fill=MID_BLUE)
    add_textbox(slide, left+0.1, 1.52, 2.8, 0.65,
                title, size=12, bold=True, color=WHITE)
    add_textbox(slide, left+0.1, 2.25, 2.8, 3.0,
                desc, size=13, color=DARK_GREY)
    if i < 3:
        add_textbox(slide, left+3.1, 3.2, 0.3, 0.4, "→",
                    size=20, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

add_rect(slide, 0.3, 5.65, 12.7, 1.55, fill=CARD_BLUE)
add_textbox(slide, 0.5, 5.72, 12.3, 0.4,
            "From Data to Training Examples", size=13, bold=True, color=DARK_BLUE)
add_textbox(slide, 0.5, 6.1, 12.3, 1.0,
            "Each training example is a complete conversation: system prompt → user query → tool call → tool result (PRE-COMPUTED) → correct assistant response. "
            "Examples cover threshold tuning, SAR backtest, clustering, segment stats, and policy Q&A — "
            "with specific gap-fixing examples written after each round of production testing.",
            size=12, color=DARK_GREY)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Fine-Tuning Process
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(slide, "The Fine-Tuning Process",
           "From general-purpose LLM to FRAML-specialised analytics assistant")
footer(slide)

ft_steps = [
    ("Base Model", "Qwen 2.5 7B Instruct\n(Unsloth 4-bit hub mirror)\nGeneral-purpose LLM\nno FRAML knowledge"),
    ("LoRA Training", "Low-Rank Adaptation\nR=16, Alpha=32\nLR=2e-4, 100 steps\nVast.ai RTX 3090 Ti\n~10 min per run"),
    ("GGUF Conversion", "Merge LoRA → base\nConvert to GGUF F16\nQuantize to Q4_K_M\n4.7 GB final size"),
    ("Deployment", "LM Studio\nlocalhost:1234/v1\nOpenAI-compatible API\n8,192 token context"),
]

for i, (title, desc) in enumerate(ft_steps):
    left = 0.3 + i * 3.2
    add_rect(slide, left, 1.45, 3.0, 3.8, fill=LIGHT_BLUE)
    add_rect(slide, left, 1.45, 3.0, 0.6, fill=MID_BLUE)
    add_textbox(slide, left+0.1, 1.52, 2.8, 0.52,
                title, size=14, bold=True, color=WHITE)
    add_textbox(slide, left+0.1, 2.1, 2.8, 3.0,
                desc, size=13, color=DARK_GREY)
    if i < 3:
        add_textbox(slide, left+3.1, 3.1, 0.3, 0.4, "→",
                    size=20, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

add_rect(slide, 0.3, 5.5, 8.4, 1.7, fill=CARD_BLUE)
add_textbox(slide, 0.5, 5.58, 8.0, 0.45,
            "Why LoRA?", size=13, bold=True, color=DARK_BLUE)
add_textbox(slide, 0.5, 6.0, 8.0, 1.1,
            "LoRA fine-tunes only ~1% of model parameters (small adapter matrices) rather than the full 7B weights. "
            "Training is fast (~10 min on a single GPU) and reversible. "
            "Base model knowledge is preserved; adapters inject FRAML-specific behaviour.",
            size=13, color=DARK_GREY)
slide.shapes.add_picture(str(IMG_LORA), Inches(8.9), Inches(5.5), Inches(4.1), Inches(1.7))

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Iterative Gap Fixing
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(slide, "Iterative Improvement: V1 → V4",
           "Production testing drives continuous gap identification and training data expansion")
footer(slide)

versions = [
    ("V1", "60 base examples\nInitial FRAML behaviour\nBasic tool calling\nEnglish-only output", "10 gaps\nidentified", False),
    ("V2", "+20 gap examples\n80 total\nFixed: language, FP/FN\nlabels, threshold\ndirection, citations", "6 gaps\nidentified", False),
    ("V3", "+36 gap examples\n116 total\nFixed: V2 gaps +\nSAR backtest tool\ncoverage", "11 gaps\nidentified", False),
    ("V4", "+47 gap examples\n163 total\nFixed: commentary\ncontrol, cluster stats,\nno fabricated thresholds,\npolicy citations", "Testing\nin progress", True),
]

for i, (ver, desc, result, is_current) in enumerate(versions):
    left = 0.3 + i * 3.2
    add_rect(slide, left, 1.45, 3.0, 4.2, fill=LIGHT_BLUE)
    banner_fill = ORANGE if is_current else MID_BLUE
    add_rect(slide, left, 1.45, 3.0, 0.65, fill=banner_fill)
    add_textbox(slide, left+0.1, 1.52, 2.8, 0.55,
                ver, size=20, bold=True, color=WHITE)
    add_textbox(slide, left+0.1, 2.2, 2.8, 2.5,
                desc, size=13, color=DARK_GREY)
    badge_fill = ORANGE if not is_current else MID_BLUE
    add_rect(slide, left+0.2, 4.6, 2.6, 0.75, fill=badge_fill)
    add_textbox(slide, left+0.25, 4.65, 2.5, 0.65,
                result, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(slide, 0.3, 5.55, 12.7, 1.65, fill=CARD_BLUE)
add_textbox(slide, 0.5, 5.62, 12.3, 0.4,
            "Gap Identification Process", size=13, bold=True, color=DARK_BLUE)
add_textbox(slide, 0.5, 6.0, 12.3, 1.1,
            "After each training run: deploy model → run standard test prompts across all tools → "
            "document failures (wrong numbers, fabricated thresholds, unsolicited commentary) → "
            "write correct gold-standard examples → retrain. "
            "Each example shows the RIGHT answer only — the model learns by imitation, not by seeing its own mistakes.",
            size=12, color=DARK_GREY)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Live Demo placeholder
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(slide, "Live Demo", None)
add_rect(slide, 0, 1.18, 13.33, 0.06, fill=ORANGE)
footer(slide)

demo_items = [
    ("Dynamic Segmentation", "Cluster Business customers into 4 groups\nIdentify highest-risk segment\nTreemap visualisation"),
    ("Threshold Tuning", "FP/FN sweep for Individual AVG_TRXN_AMT\nCrossover point analysis\nSAR backtest at 90% catch rate"),
    ("Policy Q&A", "Query structuring detection policy\nCitation from FFIEC Manual\nKB disclaimer behaviour"),
]

for i, (title, desc) in enumerate(demo_items):
    left = 0.5 + i * 4.2
    add_rect(slide, left, 1.5, 3.8, 3.5, fill=LIGHT_BLUE)
    add_rect(slide, left, 1.5, 3.8, 0.6, fill=MID_BLUE)
    add_textbox(slide, left+0.15, 1.57, 3.5, 0.5,
                title, size=14, bold=True, color=WHITE)
    add_textbox(slide, left+0.15, 2.2, 3.5, 1.35,
                desc, size=13, color=DARK_GREY)

# Mini visualisations inside demo cards
demo_imgs = [IMG_CLUSTER, IMG_FPFN, IMG_RAG]
for i, img_path in enumerate(demo_imgs):
    left = 0.5 + i * 4.2
    slide.shapes.add_picture(str(img_path), Inches(left + 0.15), Inches(3.65),
                             Inches(3.5), Inches(1.25))

add_rect(slide, 1.5, 5.3, 10.0, 1.6, fill=CARD_BLUE)
add_textbox(slide, 1.65, 5.38, 9.7, 0.5,
            "App running at http://127.0.0.1:5000",
            size=16, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
add_textbox(slide, 1.65, 5.85, 9.7, 0.9,
            "Model: qwen-framl-v5  ·  LM Studio localhost:1234/v1  ·  Context: 8,192 tokens",
            size=12, color=DARK_GREY, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# Save
# ═════════════════════════════════════════════════════════════════════════════
out = "framl_ai_presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
